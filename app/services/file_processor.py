from pathlib import Path
from typing import Union

import pdfplumber
import docx
import openpyxl
import pptx
try:
    import xlrd
except ImportError:
    xlrd = None

class FileProcessor:
    def extract_text(self, file_path: Union[str, Path]) -> str:
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return self._extract_pdf(file_path)
        elif suffix in [".docx", ".doc"]:
            return self._extract_docx(file_path)
        elif suffix == ".xlsx":
            return self._extract_xlsx(file_path)
        elif suffix == ".xls":
            return self._extract_xls(file_path)
        elif suffix in [".pptx", ".ppt"]:
            return self._extract_pptx(file_path)
        elif suffix == ".txt":
            return self._extract_txt(file_path)
        elif suffix in [".html", ".htm"]:
            return self._extract_html(file_path)
        elif suffix == ".mhtml":
            return self._extract_mhtml(file_path)
        elif suffix == ".xml":
            return self._extract_xml(file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
    def _extract_html(self, file_path: Path) -> str:
        from bs4 import BeautifulSoup
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text(separator=" ", strip=True)

    def _extract_mhtml(self, file_path: Path) -> str:
        import email
        from bs4 import BeautifulSoup
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            msg = email.message_from_file(f)
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    html = part.get_payload(decode=True)
                    if html:
                        soup = BeautifulSoup(html, "html.parser")
                        return soup.get_text(separator=" ", strip=True)
        return ""

    def _extract_xml(self, file_path: Path) -> str:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()
        return " ".join([elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()])

    def _extract_pdf(self, file_path: Path) -> str:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text

    def _extract_docx(self, file_path: Path) -> str:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])


    def _extract_xlsx(self, file_path: Path) -> str:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
        return "\n".join(text)

    def _extract_xls(self, file_path: Path) -> str:
        if xlrd is None:
            raise ImportError("xlrd is required to read .xls files. Please install it with 'poetry add xlrd'.")
        wb = xlrd.open_workbook(file_path)
        text = []
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
        return "\n".join(text)

    def _extract_pptx(self, file_path: Path) -> str:
        prs = pptx.Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_txt(self, file_path: Path) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
